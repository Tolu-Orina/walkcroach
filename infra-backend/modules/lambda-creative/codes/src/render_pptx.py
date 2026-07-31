"""WalkCroach PPTX renderer — Graphite Lumen 16:9 decks.

Coordinate system (walkcroach-pptx skill):
  Design canvas 1920×1080 → EMU @ 6350 EMU/px
  Slide size 12,192,000 × 6,858,000 EMU

HD enforcement: add_hd_image.assert_hd_fit / max_fit_box before every picture.
"""
from __future__ import annotations

import importlib.util
import sys
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

# Graphite Lumen (skills/web/walkcroach-brand-guidelines)
INK = RGBColor(0x0B, 0x0C, 0x0F)
PAPER = RGBColor(0xF2, 0xF3, 0xF5)
MIST = RGBColor(0x91, 0x98, 0xA4)
SIGNAL = RGBColor(0xF0, 0xB4, 0x29)

SLIDE_W_EMU = 12_192_000
SLIDE_H_EMU = 6_858_000
EMU_PER_PX = 6350


def px(n: float) -> Emu:
    return Emu(int(round(n * EMU_PER_PX)))


def _load_hd_helper():
    from run_skill_script import resolve_script

    path = resolve_script("add_hd_image")
    spec = importlib.util.spec_from_file_location("add_hd_image", path)
    if spec is None or spec.loader is None:
        raise RuntimeError("cannot load add_hd_image")
    mod = importlib.util.module_from_spec(spec)
    sys.modules["add_hd_image"] = mod
    spec.loader.exec_module(mod)
    return mod


def _set_run(run, text: str, *, size_pt: int, bold: bool = False, color: RGBColor = PAPER):
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _add_text(
    slide,
    text: str,
    *,
    left_px: int,
    top_px: int,
    w_px: int,
    h_px: int,
    size: int,
    bold: bool = False,
    color: RGBColor = PAPER,
):
    box = slide.shapes.add_textbox(px(left_px), px(top_px), px(w_px), px(h_px))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    _set_run(run, text, size_pt=size, bold=bold, color=color)
    return box


def _add_bullets(slide, bullets: list[str], *, left_px: int = 80, top_px: int = 200, w_px: int = 1000):
    """Real OOXML bullets — never put a literal • into a text run (validate_pptx)."""
    box = slide.shapes.add_textbox(px(left_px), px(top_px), px(w_px), px(700))
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets[:8]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(10)
        pPr = p._p.get_or_add_pPr()
        # Clear any prior bullet children, then set buFont + buChar
        for child in list(pPr):
            tag = child.tag
            if tag.endswith("}buFont") or tag.endswith("}buChar") or tag.endswith("}buNone"):
                pPr.remove(child)
        bu_font = etree.SubElement(pPr, qn("a:buFont"))
        bu_font.set("typeface", "Arial")
        bu_char = etree.SubElement(pPr, qn("a:buChar"))
        bu_char.set("char", "•")
        run = p.add_run()
        # Text only — the glyph comes from buChar, not from the run
        clean = str(bullet).lstrip("•-*–— ").strip()[:220]
        _set_run(run, clean, size_pt=18, color=PAPER)
    return box


def _fill_bg(slide, color: RGBColor = INK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_image_safe(
    slide,
    image_path: Path,
    *,
    left_px: int,
    top_px: int,
    max_w_px: int,
    max_h_px: int,
    alt_text: str = "",
):
    hd = _load_hd_helper()
    w, h = hd.max_fit_box(image_path, max_w_px, max_h_px)
    hd.assert_hd_fit(image_path, w, h)
    pic = slide.shapes.add_picture(
        str(image_path),
        px(left_px),
        px(top_px),
        width=px(w),
        height=px(h),
    )
    # Phase E3 — OOXML cNvPr descr for screen readers / check_creative_a11y
    descr = (alt_text or image_path.stem or "Slide image").strip()[:255]
    try:
        pic._element.nvPicPr.cNvPr.set("descr", descr)  # type: ignore[attr-defined]
    except Exception:
        try:
            pic.name = descr[:50]
        except Exception:
            pass
    return pic


def render_pptx(
    brief: dict[str, Any],
    out_path: Path,
    image_paths: dict[str, Path] | None = None,
) -> Path:
    """Build a 16:9 Graphite Lumen deck from a structured brief.

    brief schema:
      {
        "title": str,
        "subtitle": str?,
        "slides": [
          {"title": str, "bullets": [str], "notes": str?, "image_key": str?}
        ]
      }
    """
    image_paths = image_paths or {}
    slides_spec = brief.get("slides") or []
    if not slides_spec:
        raise ValueError("brief.slides must be a non-empty list")
    if len(slides_spec) > 12:
        raise ValueError("brief.slides capped at 12 for Phase B")

    prs = Presentation()
    prs.slide_width = SLIDE_W_EMU
    prs.slide_height = SLIDE_H_EMU
    blank = prs.slide_layouts[6]

    title_slide = prs.slides.add_slide(blank)
    _fill_bg(title_slide, INK)
    _add_text(
        title_slide,
        str(brief.get("title") or "Untitled")[:120],
        left_px=80,
        top_px=360,
        w_px=1760,
        h_px=120,
        size=48,
        bold=True,
    )
    subtitle = str(brief.get("subtitle") or "").strip()
    if subtitle:
        _add_text(
            title_slide,
            subtitle[:200],
            left_px=80,
            top_px=480,
            w_px=1400,
            h_px=80,
            size=22,
            color=MIST,
        )
    # Sparse amber CTA stripe — not a full-height accent bar
    shape = title_slide.shapes.add_shape(1, px(80), px(1000), px(160), px(6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SIGNAL
    shape.line.fill.background()

    for spec in slides_spec:
        slide = prs.slides.add_slide(blank)
        _fill_bg(slide, INK)
        _add_text(
            slide,
            str(spec.get("title") or "Slide")[:120],
            left_px=80,
            top_px=48,
            w_px=1760,
            h_px=90,
            size=32,
            bold=True,
        )
        bullets = [str(b) for b in (spec.get("bullets") or []) if str(b).strip()]
        if not bullets and spec.get("notes"):
            bullets = [str(spec["notes"])]
        if not bullets:
            bullets = ["Key point forthcoming"]

        img_key = spec.get("image_key")
        img_path = image_paths.get(str(img_key)) if img_key else None
        if img_path and Path(img_path).is_file():
            _add_bullets(slide, bullets, left_px=80, top_px=180, w_px=900)
            try:
                alt = str(
                    brief.get("altText")
                    or brief.get("alt_text")
                    or spec.get("altText")
                    or spec.get("title")
                    or "Slide illustration"
                )
                _add_image_safe(
                    slide,
                    Path(img_path),
                    left_px=1060,
                    top_px=180,
                    max_w_px=780,
                    max_h_px=720,
                    alt_text=alt,
                )
            except ValueError:
                pass  # HD refuse — skip rather than upscale
        else:
            _add_bullets(slide, bullets, left_px=80, top_px=180, w_px=1760)

        notes = str(spec.get("notes") or "").strip()
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes[:2000]

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path
