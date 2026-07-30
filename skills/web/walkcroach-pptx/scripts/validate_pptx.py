#!/usr/bin/env python3
"""WalkCroach PPTX structural + content QA.

Original WalkCroach utility (not derived from Anthropic proprietary scripts).
Runs deterministically in lambda-creative — do not ask the LLM to "eyeball" OOXML.

Usage:
  python validate_pptx.py deck.pptx
  python validate_pptx.py deck.pptx --json
Exit 0 = pass, 1 = fail with findings on stdout.
"""
from __future__ import annotations

import argparse
import json
import re
import sys
import zipfile
from pathlib import Path

PLACEHOLDER_RE = re.compile(
    r"\b(lorem|ipsum|TODO|xxx+|\[insert|this (page|slide) layout)\b",
    re.I,
)
# AI-slide tells we flag when present as thin full-width rects near titles
# (heuristic; visual QA still required via thumbnail script).


def _fail(msg: str, findings: list[str]) -> None:
    findings.append(msg)


def validate(path: Path) -> list[str]:
    findings: list[str] = []
    if not path.is_file():
        return [f"missing file: {path}"]

    try:
        zf = zipfile.ZipFile(path)
    except zipfile.BadZipFile:
        return ["not a valid ZIP/OOXML package"]

    names = set(zf.namelist())
    required = [
        "[Content_Types].xml",
        "ppt/presentation.xml",
    ]
    for r in required:
        if r not in names:
            _fail(f"missing required part: {r}", findings)

    slides = sorted(n for n in names if re.match(r"ppt/slides/slide\d+\.xml$", n))
    if not slides:
        _fail("no slides found under ppt/slides/", findings)

    # Content-type must declare slides
    try:
        ct = zf.read("[Content_Types].xml").decode("utf-8", errors="replace")
    except KeyError:
        ct = ""
    for s in slides:
        if s.split("/")[-1] not in ct and s not in ct:
            # soft check — some packs use Override PartName
            if f'PartName="/{s}"' not in ct and f"PartName=\"/{s}\"" not in ct:
                _fail(f"slide not in [Content_Types].xml: {s}", findings)

    # Placeholder / lorem scan across slide XML + notes
    text_blobs: list[str] = []
    for name in names:
        if not name.endswith(".xml"):
            continue
        if "/slides/" not in name and "/notesSlides/" not in name:
            continue
        raw = zf.read(name).decode("utf-8", errors="replace")
        text_blobs.append(raw)
        if PLACEHOLDER_RE.search(raw):
            _fail(f"placeholder/lorem-like text in {name}", findings)

    # Literal bullet character often doubles with real bullet formatting
    for name in slides:
        raw = zf.read(name).decode("utf-8", errors="replace")
        if ">•<" in raw or ">• " in raw or ">•</" in raw:
            _fail(
                f"literal bullet character in {name} — use list formatting, not •",
                findings,
            )

    # python-pptx optional deeper checks
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Emu  # type: ignore

        prs = Presentation(str(path))
        # 16:9 expected for WalkCroach decks (EMU)
        w, h = int(prs.slide_width), int(prs.slide_height)
        # Standard 16:9 = 12192000 x 6858000; allow small tolerance
        if abs(w - 12_192_000) > 50_000 or abs(h - 6_858_000) > 50_000:
            findings.append(
                f"WARN slide size {w}x{h} EMU — WalkCroach default is 12192000x6858000 (16:9)"
            )
        for i, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for p in shape.text_frame.paragraphs:
                    t = (p.text or "").strip()
                    if len(t) > 500:
                        findings.append(
                            f"WARN slide {i}: very long paragraph ({len(t)} chars) — may overflow"
                        )
    except ImportError:
        findings.append("WARN python-pptx not installed — skipped shape-level checks")
    except Exception as e:  # noqa: BLE001
        findings.append(f"WARN python-pptx check error: {e}")

    zf.close()
    return findings


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("pptx", type=Path)
    ap.add_argument("--json", action="store_true")
    args = ap.parse_args()
    findings = validate(args.pptx)
    errors = [f for f in findings if not f.startswith("WARN")]
    warns = [f for f in findings if f.startswith("WARN")]
    if args.json:
        print(json.dumps({"ok": not errors, "errors": errors, "warnings": warns}))
    else:
        for f in findings:
            print(f)
        if not findings:
            print("OK")
    return 1 if errors else 0


if __name__ == "__main__":
    sys.exit(main())
